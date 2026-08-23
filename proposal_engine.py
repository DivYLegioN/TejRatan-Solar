"""
proposal_engine.py -- Solar PPA / Open Access client proposal generator.

IMPORTANT / SINGLE SOURCE OF TRUTH:
This module performs NO electricity billing, tariff, banking, or savings
calculations of its own. It only accepts numbers that were already computed
by the existing calculator (index.html JS, mirrored by solar_core.py) and
formats them into a PPTX, a PDF and a JSON file. If a number is wrong, fix
it in the calculator -- never "adjust" it here.

Public entry point: build_proposal(payload: dict) -> dict
    payload is the JSON body posted by the frontend (see app_flask.py's
    /api/generate-proposal route). Returns a dict with proposal_id and the
    three output file paths, or raises ValidationError.
"""

from __future__ import annotations

import json
import os
import re
import shutil
import subprocess
import unicodedata
from dataclasses import dataclass, field, asdict
from datetime import datetime, date
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt


BASE_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_ROOT = os.path.join(BASE_DIR, "generated_proposals")

# ---------------------------------------------------------------------------
# Theme -- change these to re-brand every future proposal at once.
# ---------------------------------------------------------------------------
class Theme:
    NAVY = RGBColor(0x0B, 0x1F, 0x3A)
    SOLAR_ORANGE = RGBColor(0xF2, 0x8C, 0x28)
    TEAL = RGBColor(0x12, 0x7C, 0x84)
    LIGHT_BG = RGBColor(0xF7, 0xF8, 0xFA)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
    MUTED_TEXT = RGBColor(0x60, 0x6B, 0x76)
    GREEN_SAVE = RGBColor(0x1E, 0x8E, 0x3E)
    LINE = RGBColor(0xE3, 0xE6, 0xEA)

    FONT_HEAD = "Calibri"
    FONT_BODY = "Calibri"

    COMPANY_NAME_DEFAULT = "Your Solar Company"

    SLIDE_W_IN = 13.333
    SLIDE_H_IN = 7.5


class ValidationError(Exception):
    """Raised when required client/calculation data is missing or invalid."""

    def __init__(self, errors: List[str]):
        self.errors = errors
        super().__init__("; ".join(errors))


# ---------------------------------------------------------------------------
# Data model
# ---------------------------------------------------------------------------
@dataclass
class ClientInfo:
    client_name: str = ""
    company_name: str = ""
    client_type: str = ""
    location: str = ""
    contact_person: str = ""
    designation: str = ""
    email: str = ""
    phone: str = ""
    proposal_date: str = ""
    logo_path: Optional[str] = None


NA = "Not available"


def _num(d: dict, key: str, default=None):
    """Pull a numeric value from posted data. Never invents a value --
    returns None (rendered as 'Not available') if missing."""
    v = d.get(key, default)
    if v is None or v == "":
        return None
    try:
        return float(v)
    except (TypeError, ValueError):
        return None


def _txt(d: dict, key: str, default=""):
    v = d.get(key, default)
    return str(v) if v is not None else default


@dataclass
class ProposalData:
    client: ClientInfo
    electricity: Dict[str, Any] = field(default_factory=dict)   # existing situation
    ppa: Dict[str, Any] = field(default_factory=dict)            # proposed PPA / OA
    calculations: Dict[str, Any] = field(default_factory=dict)   # financial outputs
    monthly_generation: List[float] = field(default_factory=list)  # 12 values, optional
    metadata: Dict[str, Any] = field(default_factory=dict)

    def to_json_dict(self) -> dict:
        return {
            "client": asdict(self.client),
            "electricity": self.electricity,
            "ppa": self.ppa,
            "calculations": self.calculations,
            "monthly_generation": self.monthly_generation,
            "metadata": self.metadata,
        }


def parse_payload(payload: dict) -> ProposalData:
    c = payload.get("client", {}) or {}
    client = ClientInfo(
        client_name=_txt(c, "client_name"),
        company_name=_txt(c, "company_name"),
        client_type=_txt(c, "client_type"),
        location=_txt(c, "location"),
        contact_person=_txt(c, "contact_person"),
        designation=_txt(c, "designation"),
        email=_txt(c, "email"),
        phone=_txt(c, "phone"),
        proposal_date=_txt(c, "proposal_date") or date.today().strftime("%d %b %Y"),
        logo_path=c.get("logo_path"),
    )
    return ProposalData(
        client=client,
        electricity=payload.get("electricity", {}) or {},
        ppa=payload.get("ppa", {}) or {},
        calculations=payload.get("calculations", {}) or {},
        monthly_generation=payload.get("monthly_generation", []) or [],
        metadata=payload.get("metadata", {}) or {},
    )


# ---------------------------------------------------------------------------
# Validation -- required fields only. Everything else may be "Not available".
# ---------------------------------------------------------------------------
def validate(data: ProposalData) -> List[str]:
    errors = []

    if not data.client.client_name.strip():
        errors.append("Client name is required.")

    consumption = _num(data.electricity, "monthly_consumption_kwh")
    if consumption is None:
        consumption = _num(data.electricity, "annual_consumption_kwh")
    if consumption is None:
        errors.append("Monthly or annual consumption (kWh) is required.")
    elif consumption < 0:
        errors.append("Consumption cannot be negative.")

    current_tariff = _num(data.electricity, "current_tariff")
    if current_tariff is None:
        current_tariff = _num(data.electricity, "government_tariff")
    if current_tariff is None:
        errors.append("Current tariff (Rs/unit) is required.")
    elif current_tariff < 0:
        errors.append("Current tariff cannot be negative.")

    ppa_tariff = _num(data.ppa, "ppa_tariff")
    if ppa_tariff is None:
        errors.append("PPA tariff (Rs/unit) is required.")
    elif ppa_tariff < 0:
        errors.append("PPA tariff cannot be negative.")

    solar_capacity = _num(data.ppa, "solar_capacity_kw")
    if solar_capacity is None:
        solar_capacity = _num(data.ppa, "solar_capacity_mw")
    if solar_capacity is None:
        errors.append("Solar capacity is required.")
    elif solar_capacity < 0:
        errors.append("Solar capacity cannot be negative.")

    contract_period = _num(data.ppa, "contract_period_years")
    if contract_period is None:
        errors.append("Contract period (years) is required.")
    elif contract_period < 0:
        errors.append("Contract period cannot be negative.")

    # Generic negative-value sweep across every numeric field supplied.
    for section_name, section in (
        ("electricity", data.electricity),
        ("ppa", data.ppa),
        ("calculations", data.calculations),
    ):
        for k, v in section.items():
            if isinstance(v, (int, float)) and v < 0:
                errors.append(f"'{k}' in {section_name} cannot be negative (got {v}).")

    return errors


# ---------------------------------------------------------------------------
# Formatting helpers
# ---------------------------------------------------------------------------
def fmt_rs(v, decimals=0) -> str:
    if v is None:
        return NA
    try:
        v = float(v)
    except (TypeError, ValueError):
        return NA
    s = f"{v:,.{decimals}f}"
    return f"\u20b9{s}"


def fmt_num(v, decimals=1, suffix="") -> str:
    if v is None:
        return NA
    try:
        v = float(v)
    except (TypeError, ValueError):
        return NA
    return f"{v:,.{decimals}f}{suffix}"


def fmt_pct(v, decimals=1) -> str:
    if v is None:
        return NA
    try:
        v = float(v)
    except (TypeError, ValueError):
        return NA
    return f"{v:,.{decimals}f}%"


def get_any(d: dict, *keys, default=None):
    for k in keys:
        if k in d and d[k] not in (None, ""):
            return d[k]
    return default


def slugify(text: str) -> str:
    text = unicodedata.normalize("NFKD", text).encode("ascii", "ignore").decode("ascii")
    text = re.sub(r"[^\w\s-]", "", text).strip().lower()
    return re.sub(r"[-\s]+", "-", text) or "client"


# ---------------------------------------------------------------------------
# Chart generation (matplotlib) -- charts are only built from numbers that
# were actually supplied; if inputs are missing, the chart is skipped.
# ---------------------------------------------------------------------------
CHART_PALETTE = {
    "current": "#8A93A0",
    "proposed": "#F28C28",
    "save": "#127C84",
    "grid": "#E3E6EA",
}


def _style_ax(ax):
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color(CHART_PALETTE["grid"])
    ax.spines["bottom"].set_color(CHART_PALETTE["grid"])
    ax.tick_params(colors="#3A3F47", labelsize=11)
    ax.yaxis.grid(True, color=CHART_PALETTE["grid"], linewidth=1)
    ax.set_axisbelow(True)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"{v:,.0f}"))
    ax.margins(y=0.18)


def chart_annual_cost_comparison(out_path, current_annual, proposed_annual):
    if current_annual is None or proposed_annual is None:
        return None
    fig, ax = plt.subplots(figsize=(6.4, 3.6), dpi=150)
    bars = ax.bar(["Current", "Proposed (PPA)"], [current_annual, proposed_annual],
                   color=[CHART_PALETTE["current"], CHART_PALETTE["proposed"]], width=0.5)
    for b in bars:
        ax.annotate(f"\u20b9{b.get_height():,.0f}", (b.get_x() + b.get_width() / 2, b.get_height()),
                    ha="center", va="bottom", fontsize=12, fontweight="bold", color="#1A1A1A")
    ax.set_ylabel("Annual Electricity Cost (\u20b9)")
    _style_ax(ax)
    fig.tight_layout()
    fig.savefig(out_path, transparent=True)
    plt.close(fig)
    return out_path


def chart_monthly_cost_comparison(out_path, current_monthly, proposed_monthly):
    if current_monthly is None or proposed_monthly is None:
        return None
    fig, ax = plt.subplots(figsize=(6.4, 3.6), dpi=150)
    bars = ax.bar(["Current", "Proposed (PPA)"], [current_monthly, proposed_monthly],
                   color=[CHART_PALETTE["current"], CHART_PALETTE["proposed"]], width=0.5)
    for b in bars:
        ax.annotate(f"\u20b9{b.get_height():,.0f}", (b.get_x() + b.get_width() / 2, b.get_height()),
                    ha="center", va="bottom", fontsize=12, fontweight="bold", color="#1A1A1A")
    ax.set_ylabel("Monthly Electricity Cost (\u20b9)")
    _style_ax(ax)
    fig.tight_layout()
    fig.savefig(out_path, transparent=True)
    plt.close(fig)
    return out_path


def chart_cumulative_savings(out_path, annual_saving, contract_years):
    if annual_saving is None or not contract_years:
        return None
    years = list(range(1, int(contract_years) + 1))
    cumulative = [annual_saving * y for y in years]
    fig, ax = plt.subplots(figsize=(9.5, 3.8), dpi=150)
    ax.plot(years, cumulative, color=CHART_PALETTE["save"], linewidth=2.5, marker="o", markersize=4)
    ax.fill_between(years, cumulative, color=CHART_PALETTE["save"], alpha=0.12)
    last_x, last_y = years[-1], cumulative[-1]
    ax.annotate(f"\u20b9{last_y:,.0f}", (last_x, last_y), textcoords="offset points",
                xytext=(-10, 10), fontsize=12, fontweight="bold", color=CHART_PALETTE["save"])
    ax.set_xlabel("Contract Year")
    ax.set_ylabel("Cumulative Savings (\u20b9)")
    ax.set_xticks(years)
    _style_ax(ax)
    fig.tight_layout()
    fig.savefig(out_path, transparent=True)
    plt.close(fig)
    return out_path


MONTH_LABELS = ["Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec", "Jan", "Feb", "Mar"]


def chart_monthly_generation(out_path, monthly_gen: List[float]):
    if not monthly_gen or len(monthly_gen) < 2:
        return None
    labels = MONTH_LABELS[: len(monthly_gen)]
    fig, ax = plt.subplots(figsize=(9.5, 3.6), dpi=150)
    ax.bar(labels, monthly_gen, color=CHART_PALETTE["proposed"], width=0.6)
    ax.set_ylabel("Generation (units/kWh)")
    _style_ax(ax)
    fig.tight_layout()
    fig.savefig(out_path, transparent=True)
    plt.close(fig)
    return out_path


# ---------------------------------------------------------------------------
# PPTX helpers
# ---------------------------------------------------------------------------
def _blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def _rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = Theme.LINE
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _text(slide, x, y, w, h, text, size=18, bold=False, color=Theme.DARK_TEXT,
          align=PP_ALIGN.LEFT, font=None, anchor=None, line_spacing=None, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font or Theme.FONT_BODY
    return box


def _bullets(slide, x, y, w, h, items, size=15, color=Theme.DARK_TEXT, bold_lead=False,
             gap_pt=8, marker="\u2022  "):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap_pt)
        r = p.add_run()
        r.text = f"{marker}{item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = Theme.FONT_BODY
    return box


def _slide_header(slide, kicker, title, dark=False):
    bg = Theme.NAVY if dark else Theme.WHITE
    fg = Theme.WHITE if dark else Theme.DARK_TEXT
    _rect(slide, 0, 0, Theme.SLIDE_W_IN, Theme.SLIDE_H_IN, bg)
    if kicker:
        _text(slide, 0.6, 0.45, 8, 0.35, kicker.upper(), size=12, bold=True,
              color=Theme.SOLAR_ORANGE if dark else Theme.TEAL)
    _text(slide, 0.6, 0.75, 11.5, 0.7, title, size=30, bold=True, color=fg, font=Theme.FONT_HEAD)
    return fg


def _footer(slide, page_no, client_name, dark=False):
    fg = Theme.MUTED_TEXT if not dark else RGBColor(0xC9, 0xD2, 0xDC)
    _text(slide, 0.6, Theme.SLIDE_H_IN - 0.45, 6, 0.3, client_name, size=9, color=fg)
    _text(slide, Theme.SLIDE_W_IN - 1.3, Theme.SLIDE_H_IN - 0.45, 0.8, 0.3, str(page_no),
          size=9, color=fg, align=PP_ALIGN.RIGHT)


def _stat_card(slide, x, y, w, h, label, value, value_color=Theme.NAVY, value_size=26):
    _rect(slide, x, y, w, h, Theme.LIGHT_BG)
    _text(slide, x + 0.25, y + 0.18, w - 0.5, 0.35, label.upper(), size=11, bold=True, color=Theme.MUTED_TEXT)
    _text(slide, x + 0.25, y + h - 0.62, w - 0.5, 0.5, value, size=value_size, bold=True, color=value_color)


def _table(slide, x, y, w, rows: List[List[str]], col_widths=None, header=True,
           font_size=13, row_h=0.42):
    n_rows = len(rows)
    n_cols = len(rows[0])
    gshape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(row_h * n_rows))
    table = gshape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(w * cw / total)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(font_size)
            para.font.name = Theme.FONT_BODY
            if header and r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = Theme.NAVY
                para.font.color.rgb = Theme.WHITE
                para.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = Theme.WHITE if r % 2 else Theme.LIGHT_BG
                para.font.color.rgb = Theme.DARK_TEXT
    return table


def _flow_box(slide, x, y, w, h, text, fill=Theme.TEAL, text_color=Theme.WHITE, size=14):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    r.font.name = Theme.FONT_BODY
    return shp


def _arrow_down(slide, x, y, size=0.35):
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(size), Inches(size * 1.1))
    shp.fill.solid()
    shp.fill.fore_color.rgb = Theme.MUTED_TEXT
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _arrow_right(slide, x, y, w=0.5, h=0.3):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = Theme.MUTED_TEXT
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _picture_or_placeholder(slide, x, y, w, h, path):
    if path and os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        _rect(slide, x, y, w, h, Theme.LIGHT_BG)
        _text(slide, x, y + h / 2 - 0.2, w, 0.4, "Chart data not available", size=12,
              color=Theme.MUTED_TEXT, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main PPTX builder
# ---------------------------------------------------------------------------
def build_pptx(data: ProposalData, chart_paths: Dict[str, Optional[str]], out_path: str):
    prs = Presentation()
    prs.slide_width = Inches(Theme.SLIDE_W_IN)
    prs.slide_height = Inches(Theme.SLIDE_H_IN)

    cl = data.client
    elec = data.electricity
    ppa = data.ppa
    calc = data.calculations
    company_name = cl.company_name or Theme.COMPANY_NAME_DEFAULT

    current_annual = get_any(calc, "current_annual_cost")
    if current_annual is None:
        current_annual = get_any(elec, "current_annual_cost")
    proposed_annual = get_any(calc, "proposed_annual_cost")
    current_monthly = get_any(calc, "current_monthly_bill")
    if current_monthly is None:
        current_monthly = get_any(elec, "current_monthly_bill")
    proposed_monthly = get_any(calc, "proposed_monthly_cost")
    annual_saving = get_any(calc, "annual_savings")
    monthly_saving = get_any(calc, "monthly_savings")
    saving_pct = get_any(calc, "savings_percentage")
    contract_years = get_any(ppa, "contract_period_years")
    solar_capacity_kw = get_any(ppa, "solar_capacity_kw")
    solar_capacity_mw = get_any(ppa, "solar_capacity_mw")
    ppa_tariff = get_any(ppa, "ppa_tariff")
    expected_gen = get_any(ppa, "expected_annual_generation", "expected_generation_kwh")
    co2 = get_any(calc, "estimated_co2_reduction")

    n = 0

    # ---------------- 1. Cover ----------------
    s = _blank_slide(prs); n += 1
    _rect(s, 0, 0, Theme.SLIDE_W_IN, Theme.SLIDE_H_IN, Theme.NAVY)
    _rect(s, 0, Theme.SLIDE_H_IN - 0.15, Theme.SLIDE_W_IN, 0.15, Theme.SOLAR_ORANGE)
    _text(s, 0.9, 1.5, 10, 0.5, company_name.upper(), size=16, bold=True, color=Theme.SOLAR_ORANGE)
    _text(s, 0.9, 2.3, 11, 1.6, "Solar Power Purchase\nAgreement Proposal", size=42, bold=True,
          color=Theme.WHITE, font=Theme.FONT_HEAD, line_spacing=1.05)
    _text(s, 0.9, 4.3, 10, 0.5, f"Prepared for: {cl.client_name or NA}", size=18, color=Theme.WHITE)
    loc_line = cl.location or NA
    _text(s, 0.9, 4.8, 10, 0.5, loc_line, size=15, color=RGBColor(0xC9, 0xD2, 0xDC))
    _text(s, 0.9, Theme.SLIDE_H_IN - 0.9, 6, 0.4, cl.proposal_date, size=13,
          color=RGBColor(0xC9, 0xD2, 0xDC))

    # ---------------- 2. Executive Summary ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Overview", "Executive Summary")
    cards = [
        ("Current Annual Cost", fmt_rs(current_annual), Theme.NAVY),
        ("Proposed PPA Cost", fmt_rs(proposed_annual), Theme.TEAL),
        ("Annual Savings", fmt_rs(annual_saving), Theme.GREEN_SAVE),
        ("Savings %", fmt_pct(saving_pct), Theme.SOLAR_ORANGE),
    ]
    cw = 2.85
    for i, (label, value, color) in enumerate(cards):
        _stat_card(s, 0.6 + i * (cw + 0.2), 1.8, cw, 1.6, label, value, value_color=color)
    _text(s, 0.6, 3.8, 11.5, 0.4, "Contract Period", size=12, bold=True, color=Theme.MUTED_TEXT)
    _text(s, 0.6, 4.15, 11.5, 0.5,
          f"{fmt_num(contract_years, 0)} years" if contract_years else NA, size=20, bold=True, color=Theme.DARK_TEXT)
    _bullets(s, 0.6, 5.0, 11.5, 1.8, [
        "Predictable, lower per-unit power cost for the full contract term.",
        "Zero upfront capital investment required from the client.",
        "Renewable energy sourced through Open Access from a dedicated solar plant.",
    ], size=14)
    _footer(s, n, cl.client_name)

    # ---------------- 3. Client Electricity Profile ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Profile", "Client Electricity Profile")
    rows = [
        ["Parameter", "Value"],
        ["Sanctioned Load", f"{fmt_num(get_any(elec,'sanctioned_load_kw'))} kW"],
        ["Contract Demand", f"{fmt_num(get_any(elec,'contract_demand_kw'))} kW"],
        ["Peak Demand", f"{fmt_num(get_any(elec,'peak_demand_kw'))} kW"],
        ["Monthly Consumption", f"{fmt_num(get_any(elec,'monthly_consumption_kwh'))} kWh"],
        ["Annual Consumption", f"{fmt_num(get_any(elec,'annual_consumption_kwh'))} kWh"],
    ]
    _table(s, 0.6, 1.8, 6.3, rows, col_widths=[1.3, 1])
    _text(s, 7.3, 1.8, 5.4, 0.4, "Client Details", size=12, bold=True, color=Theme.MUTED_TEXT)
    detail_lines = [
        f"Client Type: {cl.client_type or NA}",
        f"Contact: {cl.contact_person or NA} ({cl.designation or NA})",
        f"Email: {cl.email or NA}",
        f"Phone: {cl.phone or NA}",
    ]
    _bullets(s, 7.3, 2.25, 5.4, 2, detail_lines, size=14, marker="")
    _footer(s, n, cl.client_name)

    # ---------------- 4. Current Electricity Cost ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Baseline", "Current Electricity Cost")
    rows = [
        ["Component", "Value"],
        ["Current Tariff", f"{fmt_num(get_any(elec,'current_tariff'), 2)} \u20b9/unit"],
        ["Demand Charges", fmt_rs(get_any(elec, "demand_charges"))],
        ["Fixed Charges", fmt_rs(get_any(elec, "fixed_charges"))],
        ["Electricity Duty", fmt_rs(get_any(elec, "electricity_duty"))],
        ["Wheeling Charges", fmt_rs(get_any(elec, "wheeling_charges"))],
        ["Other Charges", fmt_rs(get_any(elec, "other_charges"))],
        ["Current Monthly Bill", fmt_rs(current_monthly)],
        ["Current Annual Cost", fmt_rs(current_annual)],
    ]
    _table(s, 0.6, 1.8, 8.5, rows, col_widths=[1.5, 1], font_size=13.5)
    _footer(s, n, cl.client_name)

    # ---------------- 5. Current Consumption Analysis ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Baseline", "Current Consumption Analysis")
    _stat_card(s, 0.6, 1.8, 3.6, 1.4, "Monthly Consumption",
               f"{fmt_num(get_any(elec,'monthly_consumption_kwh'),0)} kWh", value_size=22)
    _stat_card(s, 4.4, 1.8, 3.6, 1.4, "Annual Consumption",
               f"{fmt_num(get_any(elec,'annual_consumption_kwh'),0)} kWh", value_size=22)
    _picture_or_placeholder(s, 0.6, 3.5, 11.5, 3.4, chart_paths.get("monthly_generation"))
    _footer(s, n, cl.client_name)

    # ---------------- 6. Proposed Solar PPA ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Proposal", "Proposed Solar PPA")
    cap_display = f"{fmt_num(solar_capacity_mw,2)} MW" if solar_capacity_mw else (
        f"{fmt_num(solar_capacity_kw,0)} kW" if solar_capacity_kw else NA)
    cards = [
        ("Solar Capacity", cap_display),
        ("PPA Tariff", f"{fmt_num(ppa_tariff,2)} \u20b9/unit" if ppa_tariff is not None else NA),
        ("Contract Duration", f"{fmt_num(contract_years,0)} yrs" if contract_years else NA),
        ("Expected Generation", f"{fmt_num(expected_gen,0)} kWh" if expected_gen else NA),
    ]
    cw = 2.85
    for i, (label, value) in enumerate(cards):
        _stat_card(s, 0.6 + i * (cw + 0.2), 2.0, cw, 1.6, label, value, value_size=20)
    _bullets(s, 0.6, 4.1, 11.5, 2.5, [
        "Dedicated solar generation capacity allocated to this consumer under Open Access.",
        f"PPA tariff fixed at {fmt_num(ppa_tariff,2)} \u20b9/unit" +
        (f", escalation: {fmt_pct(get_any(ppa,'ppa_escalation'))}" if get_any(ppa, "ppa_escalation") else ", no escalation specified") + ".",
        f"Generation degradation assumption: {fmt_pct(get_any(ppa,'generation_degradation'))}.",
    ], size=14)
    _footer(s, n, cl.client_name)

    # ---------------- 7. Open Access Structure ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "How It Works", "Open Access Structure")
    steps = ["Solar Plant", "Grid / Transmission\nNetwork", "Open Access\n(Scheduling & Accounting)", "Client"]
    box_w, box_h = 2.6, 1.0
    total_w = len(steps) * box_w + (len(steps) - 1) * 0.6
    start_x = (Theme.SLIDE_W_IN - total_w) / 2
    y = 2.6
    for i, label in enumerate(steps):
        x = start_x + i * (box_w + 0.6)
        color = Theme.TEAL if i % 2 == 0 else Theme.NAVY
        _flow_box(s, x, y, box_w, box_h, label, fill=color)
        if i < len(steps) - 1:
            _arrow_right(s, x + box_w + 0.05, y + box_h / 2 - 0.15, w=0.5, h=0.3)
    _bullets(s, 0.9, 4.3, 11.5, 2.2, [
        "Solar power generated at the plant is injected into the shared grid, not wired directly to the client.",
        "Open Access scheduling and energy accounting (state SLDC / DISCOM) attribute the corresponding units to the client.",
        "The client draws power from the grid as usual; billing reflects the solar units allocated under this PPA.",
    ], size=14)
    _footer(s, n, cl.client_name)

    # ---------------- 8. Cost Comparison ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Comparison", "Annual Cost Comparison")
    _picture_or_placeholder(s, 1.3, 1.9, 10.5, 4.7, chart_paths.get("annual_cost"))
    _footer(s, n, cl.client_name)

    # ---------------- 9. Monthly Savings ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Savings", "Monthly Savings")
    _picture_or_placeholder(s, 1.3, 1.7, 10.5, 4.0, chart_paths.get("monthly_cost"))
    _stat_card(s, 4.4, 5.8, 4.5, 1.0, "Monthly Saving", fmt_rs(monthly_saving), value_color=Theme.GREEN_SAVE)
    _footer(s, n, cl.client_name)

    # ---------------- 10. Annual Savings ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Savings", "Annual Savings")
    _stat_card(s, 1.8, 2.5, 4.5, 2.0, "Annual Saving", fmt_rs(annual_saving), value_color=Theme.GREEN_SAVE, value_size=32)
    _stat_card(s, 7.0, 2.5, 4.5, 2.0, "Savings Percentage", fmt_pct(saving_pct), value_color=Theme.SOLAR_ORANGE, value_size=32)
    _footer(s, n, cl.client_name)

    # ---------------- 11. Long-Term Savings ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Savings", "Long-Term Savings")
    _picture_or_placeholder(s, 1.0, 1.8, 11.3, 4.2, chart_paths.get("cumulative_savings"))
    total_contract_savings = get_any(calc, "total_contract_savings")
    if total_contract_savings is None and annual_saving is not None and contract_years:
        total_contract_savings = annual_saving * contract_years
    _text(s, 0.6, Theme.SLIDE_H_IN - 0.95, 11.5, 0.4,
          f"Estimated total savings over {fmt_num(contract_years,0)}-year contract: {fmt_rs(total_contract_savings)}",
          size=14, bold=True, color=Theme.NAVY)
    _footer(s, n, cl.client_name)

    # ---------------- 12. Solar Generation ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Generation", "Solar Generation")
    cuf = get_any(calc, "cuf") or get_any(ppa, "cuf")
    cards = [
        ("Expected Annual Generation", f"{fmt_num(expected_gen,0)} kWh" if expected_gen else NA),
        ("Capacity Utilisation Factor", fmt_pct(cuf) if cuf else NA),
    ]
    for i, (label, value) in enumerate(cards):
        _stat_card(s, 0.6 + i * 4.0, 1.9, 3.7, 1.4, label, value, value_size=20)
    _picture_or_placeholder(s, 0.6, 3.6, 11.5, 3.3, chart_paths.get("monthly_generation"))
    _footer(s, n, cl.client_name)

    # ---------------- 13. OA Charges Breakdown ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Charges", "Open Access Charges Breakdown")
    oa_rows = [
        ["Charge Component", "Rate (\u20b9/unit)"],
        ["PPA Tariff", fmt_num(ppa_tariff, 2)],
        ["Transmission Charges", fmt_num(get_any(ppa, "transmission_charges"), 2)],
        ["Wheeling Charges", fmt_num(get_any(ppa, "wheeling_charges"), 2)],
        ["Cross Subsidy Surcharge", fmt_num(get_any(ppa, "cross_subsidy_surcharge"), 2)],
        ["Additional Surcharge", fmt_num(get_any(ppa, "additional_surcharge"), 2)],
        ["Banking Charges", fmt_num(get_any(ppa, "banking_charges"), 2)],
        ["Losses", fmt_pct(get_any(ppa, "losses"))],
        ["Other OA Charges", fmt_num(get_any(ppa, "other_oa_charges"), 2)],
    ]
    _table(s, 0.6, 1.8, 8.0, oa_rows, col_widths=[1.6, 1], font_size=13.5)
    _footer(s, n, cl.client_name)

    # ---------------- 14. Effective Delivered Cost ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Landing Cost", "Effective Delivered Cost")
    eff_cost = get_any(calc, "effective_energy_cost")
    boxes = [f"PPA Tariff\n{fmt_num(ppa_tariff,2)} \u20b9/unit", "+", "OA Charges\n(see breakdown)", "=",
             f"Effective Cost\n{fmt_num(eff_cost,2) if eff_cost else NA} \u20b9/unit"]
    x = 0.7
    y = 3.0
    widths = [2.6, 0.5, 2.6, 0.5, 2.9]
    for i, (label, w) in enumerate(zip(boxes, widths)):
        if label in ("+", "="):
            _text(s, x, y + 0.25, w, 0.6, label, size=28, bold=True, color=Theme.MUTED_TEXT, align=PP_ALIGN.CENTER)
        else:
            fill = Theme.TEAL if "Effective" not in label else Theme.SOLAR_ORANGE
            _flow_box(s, x, y, w, 1.1, label, fill=fill, size=13)
        x += w + 0.15
    _footer(s, n, cl.client_name)

    # ---------------- 15. Environmental Impact ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Sustainability", "Environmental Impact")
    _stat_card(s, 3.0, 2.5, 7.3, 2.2, "Estimated Annual CO\u2082 Reduction",
               f"{fmt_num(co2,0)} kg" if co2 else NA, value_color=Theme.GREEN_SAVE, value_size=34)
    _footer(s, n, cl.client_name)

    # ---------------- 16. Contract Structure ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Terms", "Contract Structure")
    _bullets(s, 0.6, 1.9, 11.5, 4.5, [
        f"PPA Tenure: {fmt_num(contract_years,0)} years" if contract_years else "PPA Tenure: Not available",
        f"Billing Cycle: {get_any(ppa, 'billing_cycle') or 'As configured / monthly (default)'}",
        f"Payment Terms: {get_any(ppa, 'payment_terms') or NA}",
        f"Scheduling & Energy Accounting: {get_any(ppa, 'scheduling_terms') or 'As per applicable Open Access regulations'}",
        f"Open Access Approvals: {get_any(ppa, 'oa_responsibility') or 'To be handled by the developer, in coordination with the client'}",
    ], size=15, gap_pt=12)
    _footer(s, n, cl.client_name)

    # ---------------- 17. Implementation Timeline ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Roadmap", "Implementation Timeline")
    default_steps = ["Commercial\nDiscussion", "PPA / LOA\nSigning", "Connectivity /\nOA Process",
                      "Project\nDevelopment", "Commissioning", "Power\nSupply Starts"]
    steps = get_any(data.metadata, "timeline_steps") or default_steps
    n_steps = len(steps)
    box_w = 1.75
    gap = 0.25
    total_w = n_steps * box_w + (n_steps - 1) * gap
    start_x = (Theme.SLIDE_W_IN - total_w) / 2
    y = 2.8
    for i, label in enumerate(steps):
        x = start_x + i * (box_w + gap)
        _flow_box(s, x, y, box_w, 1.0, label, fill=Theme.NAVY if i % 2 == 0 else Theme.TEAL, size=12)
        if i < n_steps - 1:
            _arrow_right(s, x + box_w + 0.01, y + 0.35, w=gap - 0.02, h=0.25)
    _text(s, 0.6, 4.6, 11.5, 0.9,
          "Note: Timelines are indicative and depend on regulatory approvals, site readiness and Open Access processing "
          "by the relevant DISCOM/SLDC unless specific dates are agreed separately.",
          size=12, italic=True, color=Theme.MUTED_TEXT)
    _footer(s, n, cl.client_name)

    # ---------------- 18. Risk & Responsibility Matrix ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Responsibilities", "Risk & Responsibility Matrix")
    matrix_rows = [
        ["Area", "Developer", "Client", "DISCOM / SLDC / Regulator"],
        ["Plant construction & O&M", "\u2713", "-", "-"],
        ["Open Access application", "\u2713", "Supporting docs", "Approval"],
        ["Energy scheduling", "Coordinates", "-", "Executes"],
        ["Timely payment", "-", "\u2713", "-"],
        ["Grid availability", "-", "-", "\u2713"],
        ["Regulatory changes", "Monitors", "Shares impact", "Sets policy"],
    ]
    _table(s, 0.6, 1.8, 11.5, matrix_rows, col_widths=[1.4, 1, 1, 1.4], font_size=12.5)
    _footer(s, n, cl.client_name)

    # ---------------- 19. Why This PPA ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Value", "Why This PPA")
    _bullets(s, 0.6, 1.9, 11.5, 4.5, [
        "Cost predictability: a fixed/structured tariff protects against future grid tariff hikes.",
        "100% renewable power sourced without any upfront capital investment.",
        f"Estimated savings of {fmt_pct(saving_pct)} versus current electricity cost.",
        "No land, capex, or O&M burden on the client -- the developer owns and operates the plant.",
        "Supports sustainability / ESG and CO\u2082 reduction goals.",
    ], size=16, gap_pt=14)
    _footer(s, n, cl.client_name)

    # ---------------- 20. Final Proposal ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "Summary", "Final Proposal")
    rows = [
        ["Item", "Value"],
        ["Proposed PPA Tariff", f"{fmt_num(ppa_tariff,2)} \u20b9/unit"],
        ["Contract Capacity", cap_display],
        ["Contract Period", f"{fmt_num(contract_years,0)} years" if contract_years else NA],
        ["Estimated Annual Saving", fmt_rs(annual_saving)],
    ]
    _table(s, 0.6, 1.8, 7.0, rows, col_widths=[1.6, 1])
    _bullets(s, 8.0, 1.8, 4.7, 3, ["Next Steps:", "1. Confirm assumptions & sign LOI", "2. Execute PPA",
                                     "3. Initiate Open Access application"], size=14, marker="")
    _footer(s, n, cl.client_name)

    # ---------------- 21. Disclaimer ----------------
    s = _blank_slide(prs); n += 1
    _slide_header(s, "", "Disclaimer")
    _text(s, 0.6, 2.2, 11.5, 3,
          "All financial projections are indicative and depend on actual generation, applicable Open Access "
          "charges, regulatory approvals, losses, banking provisions, taxes and final PPA terms.",
          size=17, color=Theme.MUTED_TEXT, line_spacing=1.3)
    _footer(s, n, cl.client_name)

    # ---------------- 22. Contact / Thank You ----------------
    s = _blank_slide(prs); n += 1
    _rect(s, 0, 0, Theme.SLIDE_W_IN, Theme.SLIDE_H_IN, Theme.NAVY)
    _text(s, 0.9, 2.6, 10, 1, "Thank You", size=40, bold=True, color=Theme.WHITE, font=Theme.FONT_HEAD)
    _text(s, 0.9, 3.7, 10, 0.4, company_name, size=16, bold=True, color=Theme.SOLAR_ORANGE)
    contact_lines = [x for x in [cl.contact_person, cl.email, cl.phone] if x]
    _text(s, 0.9, 4.2, 10, 1, "  |  ".join(contact_lines) if contact_lines else "",
          size=13, color=RGBColor(0xC9, 0xD2, 0xDC))

    prs.save(out_path)
    return out_path


# ---------------------------------------------------------------------------
# PDF conversion
# ---------------------------------------------------------------------------
def _find_soffice() -> Optional[str]:
    for name in ("soffice", "libreoffice"):
        path = shutil.which(name)
        if path:
            return path
    return None


def convert_pptx_to_pdf(pptx_path: str, out_dir: str) -> Optional[str]:
    """Preferred path: LibreOffice headless. Returns the PDF path, or None
    if LibreOffice isn't available (caller should fall back)."""
    soffice = _find_soffice()
    if not soffice:
        return None
    try:
        subprocess.run(
            [soffice, "--headless", "--norestore", "--convert-to", "pdf", "--outdir", out_dir, pptx_path],
            check=True, timeout=120, capture_output=True,
        )
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired, FileNotFoundError):
        return None
    pdf_path = os.path.join(out_dir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
    return pdf_path if os.path.exists(pdf_path) else None


def convert_pptx_to_pdf_fallback(data: ProposalData, chart_paths: Dict[str, Optional[str]], out_path: str) -> str:
    """Reportlab fallback used ONLY when LibreOffice is unavailable on the
    deployment host. This is a simplified text/summary PDF -- not a full
    slide-for-slide rendering of the PPTX (reportlab cannot render pptx
    directly). It covers the key figures so a proposal is still deliverable."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, PageBreak)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

    cl = data.client
    elec, ppa, calc = data.electricity, data.ppa, data.calculations

    doc = SimpleDocTemplate(out_path, pagesize=A4, topMargin=2 * cm, bottomMargin=2 * cm)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("TitleBig", parent=styles["Title"], textColor=colors.HexColor("#0B1F3A"))
    h2 = ParagraphStyle("H2", parent=styles["Heading2"], textColor=colors.HexColor("#0B1F3A"), spaceBefore=14)
    body = styles["BodyText"]

    story = []
    story.append(Paragraph("Solar PPA Proposal", title_style))
    story.append(Paragraph(f"Prepared for: {cl.client_name or NA}", body))
    story.append(Paragraph(f"Location: {cl.location or NA}", body))
    story.append(Paragraph(f"Date: {cl.proposal_date}", body))
    story.append(Spacer(1, 12))

    story.append(Paragraph("Executive Summary", h2))
    exec_rows = [
        ["Current Annual Cost", fmt_rs(get_any(calc, "current_annual_cost") or get_any(elec, "current_annual_cost"))],
        ["Proposed Annual Cost", fmt_rs(get_any(calc, "proposed_annual_cost"))],
        ["Annual Savings", fmt_rs(get_any(calc, "annual_savings"))],
        ["Savings %", fmt_pct(get_any(calc, "savings_percentage"))],
        ["Contract Period", f"{fmt_num(get_any(ppa,'contract_period_years'),0)} years"],
    ]
    t = Table(exec_rows, colWidths=[8 * cm, 8 * cm])
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#F7F8FA")),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#E3E6EA")),
        ("FONTSIZE", (0, 0), (-1, -1), 10),
        ("TOPPADDING", (0, 0), (-1, -1), 6),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
    ]))
    story.append(t)
    story.append(Spacer(1, 10))

    for chart_key, title in [("annual_cost", "Annual Cost Comparison"),
                              ("cumulative_savings", "Long-Term Savings")]:
        p = chart_paths.get(chart_key)
        if p and os.path.exists(p):
            story.append(Paragraph(title, h2))
            story.append(Image(p, width=15 * cm, height=15 * cm * 0.42))

    story.append(Paragraph("Proposed Solar PPA", h2))
    ppa_rows = [
        ["Solar Capacity (kW)", fmt_num(get_any(ppa, "solar_capacity_kw"))],
        ["PPA Tariff (\u20b9/unit)", fmt_num(get_any(ppa, "ppa_tariff"), 2)],
        ["Expected Annual Generation (kWh)", fmt_num(get_any(ppa, "expected_annual_generation"))],
    ]
    t2 = Table(ppa_rows, colWidths=[8 * cm, 8 * cm])
    t2.setStyle(TableStyle([
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#E3E6EA")),
        ("FONTSIZE", (0, 0), (-1, -1), 10),
    ]))
    story.append(t2)
    story.append(Spacer(1, 14))

    story.append(Paragraph("Disclaimer", h2))
    story.append(Paragraph(
        "All financial projections are indicative and depend on actual generation, applicable Open Access "
        "charges, regulatory approvals, losses, banking provisions, taxes and final PPA terms.", body))
    story.append(Paragraph(
        "Note: this PDF was generated by the reportlab fallback because LibreOffice is not available on this "
        "server. It summarizes the key figures rather than reproducing every slide. The full 22-slide deck is "
        "available in the PPTX download.", ParagraphStyle("Note", parent=body, textColor=colors.grey, fontSize=8)))

    doc.build(story)
    return out_path


# ---------------------------------------------------------------------------
# Orchestration
# ---------------------------------------------------------------------------
def build_proposal(payload: dict) -> dict:
    data = parse_payload(payload)
    errors = validate(data)
    if errors:
        raise ValidationError(errors)

    slug = slugify(data.client.client_name)
    stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    proposal_id = f"{slug}_{stamp}"
    out_dir = os.path.join(OUTPUT_ROOT, proposal_id)
    os.makedirs(out_dir, exist_ok=True)

    # -- Charts --
    calc, ppa, elec = data.calculations, data.ppa, data.electricity
    current_annual = get_any(calc, "current_annual_cost") or get_any(elec, "current_annual_cost")
    proposed_annual = get_any(calc, "proposed_annual_cost")
    current_monthly = get_any(calc, "current_monthly_bill") or get_any(elec, "current_monthly_bill")
    proposed_monthly = get_any(calc, "proposed_monthly_cost")
    annual_saving = get_any(calc, "annual_savings")
    contract_years = get_any(ppa, "contract_period_years")

    chart_paths = {
        "annual_cost": chart_annual_cost_comparison(
            os.path.join(out_dir, "chart_annual_cost.png"), current_annual, proposed_annual),
        "monthly_cost": chart_monthly_cost_comparison(
            os.path.join(out_dir, "chart_monthly_cost.png"), current_monthly, proposed_monthly),
        "cumulative_savings": chart_cumulative_savings(
            os.path.join(out_dir, "chart_cumulative.png"), annual_saving, contract_years),
        "monthly_generation": chart_monthly_generation(
            os.path.join(out_dir, "chart_generation.png"), data.monthly_generation),
    }

    # -- PPTX --
    pptx_path = os.path.join(out_dir, "proposal.pptx")
    build_pptx(data, chart_paths, pptx_path)

    # -- PDF --
    pdf_path = os.path.join(out_dir, "proposal.pdf")
    converted = convert_pptx_to_pdf(pptx_path, out_dir)
    pdf_method = "libreoffice"
    if converted and converted != pdf_path:
        shutil.move(converted, pdf_path)
    elif not converted:
        convert_pptx_to_pdf_fallback(data, chart_paths, pdf_path)
        pdf_method = "reportlab_fallback"

    # -- JSON --
    json_path = os.path.join(out_dir, "proposal_data.json")
    json_payload = data.to_json_dict()
    json_payload["metadata"]["proposal_id"] = proposal_id
    json_payload["metadata"]["generated_at"] = datetime.now().isoformat()
    json_payload["metadata"]["pdf_method"] = pdf_method
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(json_payload, f, indent=2, ensure_ascii=False)

    return {
        "proposal_id": proposal_id,
        "pptx_path": pptx_path,
        "pdf_path": pdf_path,
        "json_path": json_path,
        "pdf_method": pdf_method,
    }
